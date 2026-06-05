"""PPT 课件生成工具 — 基于 python-pptx

输入：结构化 JSON
输出：.pptx 文件路径

JSON 格式：
{
    "title": "课程名称",
    "slides": [
        {"title": "页标题", "bullets": ["要点1", "要点2"]},
        {"title": "页标题", "code": "print('hello')", "lang": "python"},
        {"title": "页标题", "bullets": [...], "code": "..."}
    ]
}
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "data", "ppts")

# 深色主题配色
BG_COLOR = RGBColor(0x1A, 0x1A, 0x2E)       # 深蓝黑背景
TITLE_COLOR = RGBColor(0xBD, 0xBB, 0xFF)      # 淡紫标题
TEXT_COLOR = RGBColor(0xE0, 0xE0, 0xE0)       # 浅灰正文
ACCENT_COLOR = RGBColor(0x8B, 0x1F, 0xC0)     # 紫色强调


def build_pptx(data: dict) -> str:
    """根据 JSON 数据生成 PPTX 文件

    Returns:
        str: 生成的 .pptx 文件路径
    """
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    title = data.get("title", "未命名课件")
    slides_data = data.get("slides", [])

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 宽屏
    prs.slide_height = Inches(7.5)

    # ── 封面页 ──
    _add_cover_slide(prs, title, data.get("subtitle", ""))

    # ── 内容页 ──
    for i, slide in enumerate(slides_data):
        slide_title = slide.get("title", f"第{i+1}页")
        bullets = slide.get("bullets", [])
        code = slide.get("code", "")
        lang = slide.get("lang", "python")

        if code:
            _add_code_slide(prs, slide_title, bullets, code, lang)
        else:
            _add_content_slide(prs, slide_title, bullets)

    # ── 总结页 ──
    _add_end_slide(prs)

    # 保存文件
    filename = f"{title.replace('/', '_').replace('\\', '_')}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)

    return filepath


def _add_cover_slide(prs, title, subtitle=""):
    """添加封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # 标题
    left = Inches(1.5)
    top = Inches(2.2)
    width = Inches(10)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    if subtitle:
        top2 = Inches(4.0)
        txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = TEXT_COLOR
        p2.alignment = PP_ALIGN.CENTER

    # 装饰线
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(4.5), Inches(3.8), Inches(4.3), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()


def _add_content_slide(prs, title, bullets):
    """添加内容页"""
    _make_slide(prs, title)

    # 要点列表
    txBox = prs.slides[-1].shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.5), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"● {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(14)
        p.level = 0


def _add_code_slide(prs, title, bullets, code, lang):
    """添加带代码图片的页面"""
    try:
        from backend.utils.image_gen import code_to_image
        img_path = code_to_image(code, lang)
    except Exception:
        img_path = None

    _make_slide(prs, title)

    # 如果有要点，先显示要点（上半部分）
    if bullets:
        txBox = prs.slides[-1].shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11), Inches(2.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"● {bullet}"
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)

        code_top = Inches(4.2)
    else:
        code_top = Inches(2.0)

    # 插入代码图片
    if img_path and os.path.exists(img_path):
        try:
            # 获取图片实际尺寸，等比缩放
            from PIL import Image as PILImage
            with PILImage.open(img_path) as im:
                w, h = im.size
            ratio = min(11.0 / (w / 96), 3.0 / (h / 96))  # 96 DPI转英寸
            pic_w = Inches(w / 96 * ratio)
            pic_h = Inches(h / 96 * ratio)
            prs.slides[-1].shapes.add_picture(img_path, Inches(1.0), code_top, pic_w, pic_h)
        except Exception:
            _add_fallback_code(prs, code, code_top)
    else:
        _add_fallback_code(prs, code, code_top)


def _add_fallback_code(prs, code, top):
    """无图片时用文本框显示代码"""
    txBox = prs.slides[-1].shapes.add_textbox(Inches(1.0), top, Inches(11), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xA0, 0xC0, 0xE0)
    p.font.name = "Consolas"


def _make_slide(prs, title):
    """创建基础幻灯片框架"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # 标题
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # 标题下划线
    line = slide.shapes.add_shape(1, Inches(1.0), Inches(1.5), Inches(4.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()


def _add_end_slide(prs):
    """添加结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # 感谢文字
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(9), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢学习 · 智学系统"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 提示
    p2 = tf.add_paragraph()
    p2.text = "由 AI 多智能体协同生成"
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_COLOR
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
